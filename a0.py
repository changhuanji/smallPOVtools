import tkinter as tk
from tkinter import messagebox, filedialog
import os
import copy
from pptx import Presentation


# === 内部逻辑 ===

def _duplicate_slide(pres, index):
    """
    深度复制幻灯片，解决新建页面为空白的问题。
    """
    source_slide = pres.slides[index]
    slide_layout = source_slide.slide_layout
    dest_slide = pres.slides.add_slide(slide_layout)

    for shape in source_slide.shapes:
        try:
            new_el = copy.deepcopy(shape.element)
            dest_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
        except Exception as e:
            print(f"复制形状警告: {e}")

    return dest_slide


def _process_ppt_generation(template_path, txt_path, output_path, placeholder):
    try:
        # 1. 校验输入
        if not template_path or not txt_path or not output_path:
            messagebox.showwarning("提示", "请填写所有路径！")
            return

        # 2. 读取 TXT
        with open(txt_path, 'r', encoding='utf-8') as f:
            lines = [line.strip() for line in f.readlines() if line.strip()]

        if not lines:
            messagebox.showwarning("警告", "TXT 数据文件为空！")
            return

        # 3. 加载模板
        prs = Presentation(template_path)
        if len(prs.slides) == 0:
            messagebox.showerror("错误", "PPT 模板为空！")
            return

        # 4. 复制幻灯片
        target_count = len(lines)
        if target_count > 1:
            for _ in range(target_count - 1):
                _duplicate_slide(prs, 0)

        # 5. 替换文本
        for i, text_content in enumerate(lines):
            if i >= len(prs.slides): break
            slide = prs.slides[i]

            for shape in slide.shapes:
                if not shape.has_text_frame: continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if placeholder in run.text:
                            run.text = run.text.replace(placeholder, text_content)

        # 6. 保存并打开
        prs.save(output_path)
        messagebox.showinfo("成功", f"生成完毕！\n共 {len(lines)} 页。")
        os.startfile(output_path)

    except Exception as e:
        messagebox.showerror("运行错误", f"发生错误：{str(e)}")


# === 对外接口 ===
def show_ui(parent):
    top = tk.Toplevel(parent)
    top.title("PPT 批量生成器")
    top.geometry("500x400")
    top.transient(parent) # 修改点
    top.grab_set()        # 修改点

    # 居中
    top.update_idletasks()
    x = (top.winfo_screenwidth() - top.winfo_width()) // 2
    y = (top.winfo_screenheight() - top.winfo_height()) // 2
    top.geometry(f"+{x}+{y}")

    pad_opts = {'padx': 10, 'pady': 5}

    # === 辅助函数：安全地选择文件 ===
    def select_file(entry_widget, f_types):
        path = filedialog.askopenfilename(filetypes=f_types, parent=top)
        if path:
            entry_widget.delete(0, tk.END)
            entry_widget.insert(0, path)

    def select_save_file(entry_widget):
        path = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PPT", "*.pptx")], parent=top)
        if path:
            entry_widget.delete(0, tk.END)
            entry_widget.insert(0, path)

    # ================================

    # 1. 模板
    tk.Label(top, text="PPTX 模板路径:").pack(anchor="w", **pad_opts)
    entry_tmpl = tk.Entry(top)
    entry_tmpl.pack(fill="x", **pad_opts)
    tk.Button(top, text="浏览",
              command=lambda: select_file(entry_tmpl, [("PPT", "*.pptx")])).pack(anchor="e", padx=10)

    # 2. 占位符
    tk.Label(top, text="占位符 (例如 {name}):").pack(anchor="w", **pad_opts)
    entry_placeholder = tk.Entry(top)
    entry_placeholder.insert(0, "{name}")
    entry_placeholder.pack(fill="x", **pad_opts)

    # 3. Txt
    tk.Label(top, text="TXT 数据路径 (每行对应一页):").pack(anchor="w", **pad_opts)
    entry_txt = tk.Entry(top)
    entry_txt.pack(fill="x", **pad_opts)
    tk.Button(top, text="浏览",
              command=lambda: select_file(entry_txt, [("TXT", "*.txt")])).pack(anchor="e", padx=10)

    # 4. 输出
    tk.Label(top, text="输出文件路径:").pack(anchor="w", **pad_opts)
    entry_out = tk.Entry(top)
    entry_out.pack(fill="x", **pad_opts)
    tk.Button(top, text="浏览",
              command=lambda: select_save_file(entry_out)).pack(anchor="e", padx=10)

    # 执行
    def run():
        _process_ppt_generation(entry_tmpl.get(), entry_txt.get(), entry_out.get(), entry_placeholder.get())

    tk.Button(top, text="开始生成", bg="#4CAF50", fg="white", font=("Arial", 12, "bold"),
              command=run).pack(pady=20, fill="x", padx=20)