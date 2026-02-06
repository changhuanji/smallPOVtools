import tkinter as tk
from tkinter import messagebox, filedialog
import os
import platform
import subprocess
from pptx import Presentation
from pptx.dml.color import RGBColor


# === 内部逻辑 ===
def _process_modify_ppt(input_path, output_path, rgb_str, do_color, remove_spaces, remove_empty_boxes):
    try:
        if not input_path or not output_path:
            messagebox.showwarning("提示", "路径不能为空！")
            return

        # 颜色解析
        r, g, b = 0, 0, 0
        if do_color:
            try:
                rgb_clean = rgb_str.replace("，", ",").replace(" ", "")
                r, g, b = map(int, rgb_clean.split(','))
            except:
                messagebox.showerror("错误", "RGB颜色格式不正确！")
                return

        prs = Presentation(input_path)

        for slide in prs.slides:
            shapes_to_delete = []

            for shape in slide.shapes:
                if not shape.has_text_frame: continue

                # 功能C：检测空白文本框
                text_content = shape.text_frame.text.strip()
                if remove_empty_boxes:
                    if not text_content:
                        shapes_to_delete.append(shape)
                        continue

                        # 功能A & B
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        if do_color:
                            run.font.color.rgb = RGBColor(r, g, b)
                        if remove_spaces and run.text:
                            run.text = "".join(run.text.split())

            # 执行删除
            for shape in shapes_to_delete:
                sp = shape._element
                sp.getparent().remove(sp)

        prs.save(output_path)

        # 自动打开
        if platform.system() == 'Windows':
            os.startfile(output_path)
        else:
            subprocess.call(('open', output_path))

    except Exception as e:
        messagebox.showerror("错误", f"处理失败：{str(e)}")


# === 对外接口 ===
def show_ui(parent):
    top = tk.Toplevel(parent)
    top.title("PPT 改色与清理工具")
    top.geometry("500x450")
    top.transient(parent) # 修改点
    top.grab_set()        # 修改点

    top.update_idletasks()
    x = (top.winfo_screenwidth() - top.winfo_width()) // 2
    y = (top.winfo_screenheight() - top.winfo_height()) // 2
    top.geometry(f"+{x}+{y}")

    pad_opts = {'padx': 10, 'pady': 5}

    # === 辅助函数 ===
    def select_file(entry_widget):
        path = filedialog.askopenfilename(filetypes=[("PPT", "*.pptx")], parent=top)
        if path:
            entry_widget.delete(0, tk.END)
            entry_widget.insert(0, path)

    def select_save(entry_widget):
        path = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PPT", "*.pptx")], parent=top)
        if path:
            entry_widget.delete(0, tk.END)
            entry_widget.insert(0, path)

    # ==============

    # 1. 输入
    tk.Label(top, text="选择 PPT 文件:").pack(anchor="w", **pad_opts)
    entry_in = tk.Entry(top)
    entry_in.pack(fill="x", **pad_opts)
    tk.Button(top, text="浏览",
              command=lambda: select_file(entry_in)).pack(anchor="e", padx=10)

    # 分隔
    tk.Frame(top, height=2, bd=1, relief="sunken").pack(fill="x", padx=10, pady=10)

    # 2. 颜色
    frame_color = tk.Frame(top)
    frame_color.pack(fill="x", **pad_opts)

    var_do_color = tk.BooleanVar(value=False)
    chk_color = tk.Checkbutton(frame_color, text="修改字体颜色", variable=var_do_color)
    chk_color.pack(side="left")

    tk.Label(frame_color, text=" RGB:").pack(side="left")
    entry_rgb = tk.Entry(frame_color, width=15);
    entry_rgb.insert(0, "0,0,0");
    entry_rgb.pack(side="left")

    # 3. 清理
    var_space = tk.BooleanVar(value=False)
    tk.Checkbutton(top, text="去除所有空格/换行符", variable=var_space).pack(anchor="w", padx=10)

    var_empty_box = tk.BooleanVar(value=True)
    tk.Checkbutton(top, text="删除所有空白文本框", variable=var_empty_box).pack(anchor="w", padx=10)

    # 分隔
    tk.Frame(top, height=2, bd=1, relief="sunken").pack(fill="x", padx=10, pady=10)

    # 4. 输出
    tk.Label(top, text="保存路径:").pack(anchor="w", **pad_opts)
    entry_out = tk.Entry(top)
    entry_out.pack(fill="x", **pad_opts)
    tk.Button(top, text="浏览",
              command=lambda: select_save(entry_out)).pack(anchor="e", padx=10)

    def run():
        _process_modify_ppt(entry_in.get(), entry_out.get(), entry_rgb.get(),
                            var_do_color.get(), var_space.get(), var_empty_box.get())

    tk.Button(top, text="执行并打开", bg="#2196F3", fg="white", font=("Arial", 12, "bold"),
              command=run).pack(pady=15, fill="x", padx=20)